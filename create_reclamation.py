# create_reclamation.py
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from datetime import datetime
import os
import shutil
import config

class CreateReclamationWindow:
    """NCR window to create and modify reclamation"""
    def __init__(self, root, db, rec_id=None):
        self.root = root
        self.db = db
        self.rec_id = rec_id
        self.edit_mode = rec_id is not None
        
        self.root.title("Редактирование NCR" if self.edit_mode else "Добавление нового NCR")
        self.root.state('zoomed')

        # initialize STATIC DATA before widget creation
        self.models = []
        self.suppliers = []
        self.commodities = []
        self.parts_disposal = []
        self.repetition_list = []
        self.defects_list = []
        
        # Словарь для хранения комбобоксов
        self.comboboxes = {}
        
        # Список загруженных файлов
        self.attached_files = []
        self.file_paths = {}  # {имя_файла: полный_путь}
        self.thumbnail_images = []  # Для хранения ссылок на миниатюры
        self.thumbnails = {}  # {имя_файла: PhotoImage}
        self.selected_file = None
        
        # Стили
        self.setup_styles()
        
        # Словарь для хранения переменных полей
        self.vars = {}
        
        # Создаем status_var
        self.status_var = tk.StringVar(value="⏳ Загрузка данных...")
        
        # ✅ Сначала ЗАГРУЖАЕМ ДАННЫЕ (до создания виджетов)
        self.load_reference_data()
        if self.edit_mode:
            self.load_reclamation_data()
        
        # ✅ ПОТОМ создаем виджеты с уже загруженными данными
        self.create_widgets()
        
        # Если редактирование — загружаем файлы (после создания виджетов)
        if self.edit_mode:
            self.load_attached_files()
        
        # Настройка горячих клавиш
        self.setup_bindings()
        
        # ✅ Устанавливаем фокус
        self.root.focus_force()
    
    def setup_styles(self):
        """Настройка стилей"""
        style = ttk.Style()
        style.configure("Title.TLabel", font=('Arial', 18, 'bold'))
        style.configure("Field.TLabel", font=('Arial', 10))
        style.configure("Accent.TButton", font=('Arial', 10, 'bold'))
        style.configure("Status.TLabel", font=('Arial', 9))
        style.configure("Info.TLabel", font=('Arial', 9), foreground='gray')
    
    def create_widgets(self):
        # Главный контейнер
        main_frame = ttk.Frame(self.root, padding=20)
        main_frame.pack(fill='both', expand=True)
        
        # Верхняя панель
        top_frame = ttk.Frame(main_frame)
        top_frame.pack(fill='x', pady=(0, 20))
        
        title = ttk.Label(top_frame, text="New Non Conformity Report", 
                          style="Title.TLabel")
        title.pack(side='left')
        
        # Кнопки в верхней панели
        btn_frame = ttk.Frame(top_frame)
        btn_frame.pack(side='right')
        
        ttk.Button(btn_frame, text="🔄 Обновить справочники", command=self.refresh_reference_data).pack(side='left', padx=5)
        
        # ✅ Добавляем кнопку экспорта в PowerPoint
        ttk.Button(btn_frame, text="📊 Экспорт в PPT", command=self.export_to_pptx).pack(side='left', padx=5)
        
        # Контейнер для полей ввода (две колонки)
        fields_container = ttk.Frame(main_frame)
        fields_container.pack(fill='both', expand=True)
        
        # Левая колонка (все поля ввода)
        left_frame = ttk.LabelFrame(fields_container, text="Информация о рекламации", padding=15)
        left_frame.pack(side='left', fill='both', expand=True, padx=(0, 10))
        
        # Правая колонка (галерея изображений)
        right_frame = ttk.LabelFrame(fields_container, text="Галерея изображений", padding=15)
        right_frame.pack(side='right', fill='both', expand=True, padx=(10, 0))
        
        # ==================== ЛЕВАЯ КОЛОНКА ====================
        
        # ✅ Определяем все поля в левой колонке
        left_fields = [
            ('model', 'Модель *', 'combobox', self.models),
            ('commodity', 'Товарная группа', 'combobox', self.commodities),
            ('creator', 'Создатель', 'entry', None),
            ('partnumber', 'Номер детали *', 'entry', None),
            ('partname', 'Наименование детали', 'entry', None),
            ('supplier', 'Поставщик *', 'combobox', self.suppliers),
            ('full_pir_number', 'Номер PIR', 'entry', None),
            ('repetition', 'Повторяемость', 'combobox', self.repetition_list),
            ('ncr_status', 'Статус NCR', 'combobox', config.NCR_STATUSES),
            ('failure_quantity', 'Количество отказов', 'entry', None),
            ('vin', 'VIN номер', 'entry', None),
            ('defect', 'Дефект', 'combobox', self.defects_list),
            ('parts_disposal', 'Утилизация деталей', 'combobox', self.parts_disposal),
        ]
        
        # Создаем поля на левой панели
        for key, label, widget_type, values in left_fields:
            self._create_field(left_frame, key, label, widget_type, values)
        
        # Поле Cutoff (выбор даты)
        cutoff_frame = ttk.Frame(left_frame)
        cutoff_frame.pack(fill='x', pady=5)
        
        ttk.Label(cutoff_frame, text="Cutoff дата:", style="Field.TLabel", 
                  width=20, anchor='e').pack(side='left', padx=(0, 10))
        
        cutoff_value = self.vars.get('cutoff', tk.StringVar()).get() if 'cutoff' in self.vars else ''
        self.vars['cutoff'] = tk.StringVar(value=cutoff_value)
        self.cutoff_entry = ttk.Entry(cutoff_frame, textvariable=self.vars['cutoff'], width=30)
        self.cutoff_entry.pack(side='left', fill='x', expand=True)
        
        ttk.Button(cutoff_frame, text="📅 Выбрать дату", 
                  command=self.choose_cutoff_date).pack(side='left', padx=5)
        
        # Чекбокс Report 8D
        checkbox_frame = ttk.Frame(left_frame)
        checkbox_frame.pack(fill='x', pady=5)
        
        ttk.Label(checkbox_frame, text="", width=20, anchor='e').pack(side='left', padx=(0, 10))
        
        checkbox_value = self.vars.get('report8d_checkbox', tk.BooleanVar()).get() if 'report8d_checkbox' in self.vars else False
        self.vars['report8d_checkbox'] = tk.BooleanVar(value=checkbox_value)
        
        self.report8d_checkbox = ttk.Checkbutton(
            checkbox_frame,
            text="✅ Report 8D требуется",
            variable=self.vars['report8d_checkbox']
        )
        self.report8d_checkbox.pack(side='left', anchor='w')
        
        # Поле для описания (многострочное)
        desc_frame = ttk.LabelFrame(left_frame, text="Описание проблемы", padding=10)
        desc_frame.pack(fill='both', expand=True, pady=(10, 0))
        
        self.description_text = tk.Text(desc_frame, height=6, font=('Arial', 10), wrap='word')
        self.description_text.pack(fill='both', expand=True)
        
        if hasattr(self, '_description_value') and self._description_value:
            self.description_text.insert('1.0', self._description_value)
        
        # Поле для комментариев (многострочное)
        comments_frame = ttk.LabelFrame(left_frame, text="Комментарии", padding=10)
        comments_frame.pack(fill='both', expand=True, pady=(10, 0))
        
        self.comments_text = tk.Text(comments_frame, height=3, font=('Arial', 10), wrap='word')
        self.comments_text.pack(fill='both', expand=True)
        
        if hasattr(self, '_comments_value') and self._comments_value:
            self.comments_text.insert('1.0', self._comments_value)
        
        # ==================== ПРАВАЯ КОЛОНКА (Галерея) ====================
        
        # Кнопки управления файлами
        file_btn_frame = ttk.Frame(right_frame)
        file_btn_frame.pack(fill='x', pady=(0, 10))
        
        ttk.Button(file_btn_frame, text="📎 Добавить файлы", 
                  command=self.add_file).pack(side='left', padx=5)
        
        ttk.Button(file_btn_frame, text="🗑️ Удалить выбранный", 
                  command=self.remove_selected_file).pack(side='left', padx=5)
        
        ttk.Button(file_btn_frame, text="📂 Открыть папку", 
                  command=self.open_files_folder).pack(side='left', padx=5)
        
        # Список файлов (слева)
        list_frame = ttk.Frame(right_frame)
        list_frame.pack(side='left', fill='y', padx=(0, 10))
        
        ttk.Label(list_frame, text="Файлы:", font=('Arial', 9)).pack(anchor='w')
        
        self.file_listbox = tk.Listbox(list_frame, height=12, width=20)
        self.file_listbox.pack(side='left', fill='both', expand=True)
        self.file_listbox.bind('<<ListboxSelect>>', self.on_file_select)
        self.file_listbox.bind('<Double-Button-1>', self.on_file_double_click)
        
        file_scrollbar = ttk.Scrollbar(list_frame, orient='vertical', command=self.file_listbox.yview)
        file_scrollbar.pack(side='right', fill='y')
        self.file_listbox.configure(yscrollcommand=file_scrollbar.set)
        
        # Галерея изображений (справа)
        gallery_frame = ttk.Frame(right_frame)
        gallery_frame.pack(side='left', fill='both', expand=True)
        
        canvas_frame = ttk.Frame(gallery_frame)
        canvas_frame.pack(fill='both', expand=True)
        
        self.canvas = tk.Canvas(canvas_frame, bg='#f0f0f0', highlightthickness=0)
        scrollbar = ttk.Scrollbar(canvas_frame, orient='vertical', command=self.canvas.yview)
        
        self.canvas.configure(yscrollcommand=scrollbar.set)
        self.canvas.pack(side='left', fill='both', expand=True)
        scrollbar.pack(side='right', fill='y')
        
        self.gallery_frame = ttk.Frame(self.canvas)
        self.canvas.create_window((0, 0), window=self.gallery_frame, anchor='nw')
        
        self.gallery_frame.bind('<Configure>', self._on_gallery_configure)
        
        # Нижняя панель с кнопкой Сохранить
        bottom_btn_frame = ttk.Frame(main_frame)
        bottom_btn_frame.pack(fill='x', pady=(10, 0))
        
        ttk.Button(bottom_btn_frame, text="💾 Сохранить (Ctrl+S)", 
                  command=self.save_reclamation,
                  style="Accent.TButton").pack(side='left', padx=5)
    
    def _on_gallery_configure(self, event):
        """Обновляет область прокрутки при изменении размера галереи"""
        self.canvas.configure(scrollregion=self.canvas.bbox('all'))
    
    def _create_field(self, parent, key, label, widget_type, values):
        """Создает одно поле ввода с уже загруженными данными"""
        frame = ttk.Frame(parent)
        frame.pack(fill='x', pady=5)
        
        label_text = label + (' *' if '*' in label else '')
        ttk.Label(frame, text=label_text + ':', style="Field.TLabel", 
                width=20, anchor='e').pack(side='left', padx=(0, 10))
        
        if widget_type == 'entry':
            default_value = self.vars.get(key, tk.StringVar()).get() if key in self.vars else ''
            self.vars[key] = tk.StringVar(value=default_value)
            widget = ttk.Entry(frame, textvariable=self.vars[key], width=30)
            widget.pack(side='left', fill='x', expand=True)
            
            # ✅ Валидация для VIN
            if key == 'vin':
                vcmd = (self.root.register(self._validate_vin), '%P')
                widget.config(validate='key', validatecommand=vcmd)
                widget.bind('<FocusIn>', lambda e: self.status_var.set("Введите VIN (17 символов, буквы и цифры)"))
            else:
                widget.bind('<FocusIn>', lambda e, k=label: self.status_var.set(f"Введите {k.lower()}"))
            
            widget.bind('<FocusOut>', lambda e: self.status_var.set("Готов к работе"))
        
        elif widget_type == 'combobox':
            default_value = self.vars.get(key, tk.StringVar()).get() if key in self.vars else ''
            self.vars[key] = tk.StringVar(value=default_value)
            widget = ttk.Combobox(frame, textvariable=self.vars[key], 
                                values=values if values else [], width=28)
            widget.pack(side='left', fill='x', expand=True)
            if not default_value:
                widget.set('')
            
            self.comboboxes[key] = widget
            
            widget.bind('<FocusIn>', lambda e, k=label: self.status_var.set(f"Введите или выберите {k.lower()}"))
            widget.bind('<FocusOut>', lambda e: self.status_var.set("Готов к работе"))
            widget.bind('<KeyRelease>', self._on_combobox_keyrelease)
        
        return widget
    
    def _on_combobox_keyrelease(self, event):
        """Обработка ввода в комбобоксе для автодополнения"""
        widget = event.widget
        current_text = widget.get()
        
        if len(current_text) < 2:
            return
        
        if 'model' in str(widget):
            matches = [m for m in self.models if current_text.lower() in m.lower()]
            if matches:
                widget['values'] = matches
        elif 'supplier' in str(widget):
            matches = [s for s in self.suppliers if current_text.lower() in s.lower()]
            if matches:
                widget['values'] = matches
        elif 'commodity' in str(widget):
            matches = [c for c in self.commodities if current_text.lower() in c.lower()]
            if matches:
                widget['values'] = matches
        elif 'parts_disposal' in str(widget):
            matches = [p for p in self.parts_disposal if current_text.lower() in p.lower()]
            if matches:
                widget['values'] = matches
        elif 'repetition' in str(widget):
            matches = [r for r in self.repetition_list if current_text.lower() in r.lower()]
            if matches:
                widget['values'] = matches
        elif 'defect' in str(widget):
            matches = [d for d in self.defects_list if current_text.lower() in d.lower()]
            if matches:
                widget['values'] = matches

    def _validate_vin(self, text):
        """Ограничивает ввод VIN до 17 символов (буквы и цифры)"""
        return len(text) <= 17 and (not text or text.isalnum())
    
    def choose_cutoff_date(self):
        """Открывает диалог выбора даты для cutoff"""
        try:
            from tkcalendar import DateEntry
            
            date_window = tk.Toplevel(self.root)
            date_window.title("Выбор даты cutoff")
            date_window.geometry("320x280")
            date_window.transient(self.root)
            date_window.grab_set()
            date_window.focus_force()
            
            date_window.update_idletasks()
            width = date_window.winfo_width()
            height = date_window.winfo_height()
            x = (date_window.winfo_screenwidth() // 2) - (width // 2)
            y = (date_window.winfo_screenheight() // 2) - (height // 2)
            date_window.geometry(f'{width}x{height}+{x}+{y}')
            
            ttk.Label(date_window, text="Выберите дату cutoff:", 
                    font=('Arial', 12)).pack(pady=20)
            
            cal = DateEntry(
                date_window, 
                width=12, 
                background='darkblue',
                foreground='white',
                borderwidth=2,
                date_pattern='yyyy-mm-dd'
            )
            cal.pack(pady=20)
            cal.focus_force()
            
            result = {'date': None}
            
            def set_date():
                try:
                    date_str = cal.get()
                    if date_str:
                        result['date'] = date_str
                        date_window.destroy()
                    else:
                        messagebox.showwarning("Внимание", "Пожалуйста, выберите дату", parent=self.root)
                except Exception as e:
                    messagebox.showerror("Ошибка", f"Не удалось получить дату: {e}", parent=self.root)
            
            def clear_date():
                result['date'] = ''
                date_window.destroy()
            
            def cancel():
                result['date'] = None
                date_window.destroy()
            
            btn_frame = ttk.Frame(date_window)
            btn_frame.pack(pady=10)
            
            ttk.Button(btn_frame, text="✅ Выбрать", command=set_date, width=12).pack(side='left', padx=5)
            ttk.Button(btn_frame, text="🗑️ Очистить", command=clear_date, width=12).pack(side='left', padx=5)
            ttk.Button(btn_frame, text="❌ Отмена", command=cancel, width=12).pack(side='left', padx=5)
            
            date_window.protocol("WM_DELETE_WINDOW", cancel)
            self.root.wait_window(date_window)
            
            if result['date'] is not None:
                if result['date'] == '':
                    self.vars['cutoff'].set('')
                    self.status_var.set("🗑️ Дата cutoff очищена")
                else:
                    self.vars['cutoff'].set(result['date'])
                    self.status_var.set(f"✅ Дата cutoff выбрана: {result['date']}")
                
        except ImportError:
            messagebox.showinfo(
                "Информация", 
                "Для удобного выбора даты установите библиотеку tkcalendar:\n"
                "pip install tkcalendar\n\n"
                "Или введите дату вручную в формате ГГГГ-ММ-ДД",
                parent=self.root
            )
        except Exception as e:
            messagebox.showerror("Ошибка", f"Ошибка при выборе даты:\n{str(e)}", parent=self.root)
    
    def load_reference_data(self):
        """Загружает справочные данные из базы"""
        try:
            self.models = self.db.get_models()
            self.suppliers = self.db.get_suppliers()
            self.commodities = self.db.get_commodities()
            self.parts_disposal = self.db.get_parts_disposal()
            self.repetition_list = self.db.get_repetition_list()
            self.defects_list = self.db.get_defect_names()
            
            self.update_comboboxes()
            self.status_var.set(f"✅ Загружено: {len(self.models)} моделей, {len(self.suppliers)} поставщиков, {len(self.commodities)} товарных групп, {len(self.parts_disposal)} вариантов утилизации, {len(self.repetition_list)} вариантов повторяемости, {len(self.defects_list)} дефектов")
            
        except Exception as e:
            print(f"❌ Ошибка загрузки справочных данных: {e}")
            self.models = config.MODELS
            self.suppliers = config.SUPPLIERS
            self.commodities = []
            self.parts_disposal = []
            self.repetition_list = []
            self.defects_list = []
            self.update_comboboxes()
            self.status_var.set("⚠️ Используются значения по умолчанию")
    
    def load_reclamation_data(self):
        """Загружает данные рекламации для редактирования ДО создания виджетов"""
        self.status_var.set(f"⏳ Загрузка рекламации #{self.rec_id}...")
        
        rec = self.db.get_reclamation_by_id(self.rec_id)
        
        if rec:
            self.vars['model'] = tk.StringVar(value=rec.model or '')
            self.vars['commodity'] = tk.StringVar(value=rec.commodity or '')
            self.vars['creator'] = tk.StringVar(value=rec.creator or '')
            self.vars['partnumber'] = tk.StringVar(value=rec.partnumber or '')
            self.vars['partname'] = tk.StringVar(value=rec.partname or '')
            self.vars['supplier'] = tk.StringVar(value=rec.supplier or '')
            self.vars['full_pir_number'] = tk.StringVar(value=rec.full_pir_number or '')
            self.vars['repetition'] = tk.StringVar(value=rec.repetition or '')
            self.vars['ncr_status'] = tk.StringVar(value=rec.ncr_status or '')
            self.vars['failure_quantity'] = tk.StringVar(value=str(rec.failure_quantity) if rec.failure_quantity else '')
            self.vars['vin'] = tk.StringVar(value=rec.vin or '')
            self.vars['defect'] = tk.StringVar(value=rec.defect or '')
            self.vars['parts_disposal'] = tk.StringVar(value=rec.parts_disposal or '')
            self.vars['cutoff'] = tk.StringVar(value=rec.cutoff.strftime('%Y-%m-%d') if rec.cutoff else '')
            self.vars['report8d_checkbox'] = tk.BooleanVar(value=rec.report8d_checkbox or False)
            
            self._description_value = rec.description or ''
            self._comments_value = rec.comments or ''
            
            self.status_var.set(f"✅ Рекламация #{self.rec_id} загружена для редактирования")
        else:
            messagebox.showerror("Ошибка", f"Рекламация #{self.rec_id} не найдена", parent=self.root)
            self.status_var.set("❌ Ошибка загрузки")
    
    def load_attached_files(self):
        """Загружает список прикрепленных файлов для существующей рекламации"""
        if not self.rec_id:
            return
        
        folder_path = self.get_ncr_folder()
        if os.path.exists(folder_path):
            try:
                files = os.listdir(folder_path)
                for file in files:
                    full_path = os.path.join(folder_path, file)
                    if os.path.isfile(full_path):
                        self.attached_files.append(file)
                        self.file_paths[file] = full_path
                        self.file_listbox.insert(tk.END, file)
                self.render_gallery()
            except Exception as e:
                print(f"Ошибка загрузки файлов: {e}")
    
    def get_ncr_folder(self):
        """Возвращает путь к папке для файлов рекламации"""
        base_path = r"X:\SQE_DB\NCR"
        if self.rec_id:
            return os.path.join(base_path, str(self.rec_id))
        return base_path
    
    def add_file(self):
        """Добавляет файл к рекламации и копирует его в папку"""
        file_paths = filedialog.askopenfilenames(
            title="Выберите файлы для прикрепления",
            filetypes=[
                ("Все файлы", "*.*"),
                ("Документы", "*.pdf;*.doc;*.docx;*.xls;*.xlsx"),
                ("Изображения", "*.jpg;*.jpeg;*.png;*.bmp;*.gif"),
                ("Архивы", "*.zip;*.rar;*.7z")
            ]
        )
        
        if not file_paths:
            return
        
        for file_path in file_paths:
            file_name = os.path.basename(file_path)
            
            if file_name in self.attached_files:
                messagebox.showwarning("Внимание", f"Файл '{file_name}' уже добавлен", parent=self.root)
                continue
            
            self.attached_files.append(file_name)
            
            if self.rec_id:
                dest_path = self._copy_file_to_ncr_folder(file_path, file_name)
                if dest_path:
                    self.file_paths[file_name] = dest_path
                else:
                    self.file_paths[file_name] = file_path
            else:
                self.file_paths[file_name] = file_path
            
            self.file_listbox.insert(tk.END, file_name)
        
        self.render_gallery()
        self.status_var.set(f"✅ Добавлено файлов: {len(file_paths)}")
    
    def _copy_file_to_ncr_folder(self, source_path, file_name):
        """Копирует файл в папку рекламации и возвращает новый путь"""
        if not self.rec_id:
            return None
        
        folder_path = self.get_ncr_folder()
        if not os.path.exists(folder_path):
            os.makedirs(folder_path)
        
        dest_path = os.path.join(folder_path, file_name)
        
        if os.path.normpath(source_path) == os.path.normpath(dest_path):
            return dest_path
        
        if os.path.exists(dest_path):
            return dest_path
        
        try:
            shutil.copy2(source_path, dest_path)
            print(f"✅ Файл скопирован: {dest_path}")
            return dest_path
        except Exception as e:
            print(f"❌ Ошибка копирования: {e}")
            messagebox.showerror("Ошибка", f"Не удалось скопировать файл:\n{e}", parent=self.root)
            return None
    
    def save_files_after_save(self):
        """Обновляет пути к файлам после сохранения рекламации"""
        if not self.rec_id or not self.attached_files:
            return
        
        folder_path = self.get_ncr_folder()
        if not os.path.exists(folder_path):
            os.makedirs(folder_path)
        
        for file_name in list(self.attached_files):
            source_path = self.file_paths.get(file_name)
            if not source_path or not os.path.exists(source_path):
                continue
            
            dest_path = os.path.join(folder_path, file_name)
            
            if os.path.normpath(source_path) != os.path.normpath(dest_path):
                if not os.path.exists(dest_path):
                    try:
                        shutil.copy2(source_path, dest_path)
                        self.file_paths[file_name] = dest_path
                        print(f"✅ Файл скопирован при сохранении: {dest_path}")
                    except Exception as e:
                        print(f"❌ Ошибка копирования файла при сохранении: {e}")
                else:
                    self.file_paths[file_name] = dest_path
                    print(f"⏭️ Файл уже существует: {dest_path}")
    
    def render_gallery(self):
        """Отображает все изображения в виде галереи миниатюр"""
        for widget in self.gallery_frame.winfo_children():
            widget.destroy()
        self.thumbnail_images.clear()
        self.thumbnails = {}
        
        image_extensions = ('.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff', '.webp')
        sorted_files = sorted(self.attached_files)
        
        for file_name in sorted_files:
            file_path = self.file_paths.get(file_name)
            if not file_path or not os.path.exists(file_path):
                continue
            
            if file_name.lower().endswith(image_extensions):
                try:
                    self.create_thumbnail(file_name, file_path)
                except Exception as e:
                    print(f"Ошибка создания миниатюры для {file_name}: {e}")
        
        self.canvas.configure(scrollregion=self.canvas.bbox('all'))
    
    def create_thumbnail(self, file_name, file_path):
        """Создает миниатюру для изображения"""
        try:
            from PIL import Image, ImageTk
            
            img = Image.open(file_path)
            thumb_size = (120, 120)
            img.thumbnail(thumb_size, Image.Resampling.LANCZOS)
            
            photo = ImageTk.PhotoImage(img)
            self.thumbnails[file_name] = photo
            self.thumbnail_images.append(photo)
            
            thumb_frame = ttk.Frame(self.gallery_frame)
            thumb_frame.pack(side='left', padx=5, pady=5)
            
            img_label = ttk.Label(thumb_frame, image=photo, relief='solid', cursor='hand2')
            img_label.pack()
            
            name_label = ttk.Label(thumb_frame, text=file_name[:15] + ('...' if len(file_name) > 15 else ''), 
                                  font=('Arial', 7), foreground='gray')
            name_label.pack()
            
            img_label.bind('<Button-1>', lambda e, f=file_name: self.on_thumbnail_click(f))
            img_label.bind('<Double-Button-1>', lambda e, f=file_path: self.open_file(f))
            img_label.bind('<Button-3>', lambda e, f=file_name: self.show_thumbnail_context_menu(e, f))
            
        except Exception as e:
            print(f"Ошибка создания миниатюры: {e}")
    
    def on_thumbnail_click(self, file_name):
        """Обработка клика по миниатюре"""
        try:
            index = self.attached_files.index(file_name)
            self.file_listbox.selection_clear(0, tk.END)
            self.file_listbox.selection_set(index)
            self.file_listbox.see(index)
            self.selected_file = file_name
            self.status_var.set(f"📷 Выбрано изображение: {file_name}")
        except ValueError:
            pass
    
    def show_thumbnail_context_menu(self, event, file_name):
        """Контекстное меню для миниатюры"""
        menu = tk.Menu(self.root, tearoff=0)
        menu.add_command(label="📂 Открыть", command=lambda: self.open_file(self.file_paths.get(file_name)))
        menu.add_command(label="🗑️ Удалить", command=lambda: self.remove_file(file_name))
        menu.post(event.x_root, event.y_root)
    
    def open_file(self, file_path):
        """Открывает файл в приложении по умолчанию"""
        if not file_path or not os.path.exists(file_path):
            messagebox.showerror("Ошибка", f"Файл не найден:\n{file_path}", parent=self.root)
            return
        
        try:
            os.startfile(file_path)
            self.status_var.set(f"📂 Открыт файл: {os.path.basename(file_path)}")
        except Exception as e:
            messagebox.showerror("Ошибка", f"Не удалось открыть файл:\n{e}", parent=self.root)
    
    def remove_file(self, file_name):
        """Удаляет файл по имени"""
        if file_name not in self.attached_files:
            return
        
        if not messagebox.askyesno("Подтверждение", f"Удалить файл '{file_name}'?", parent=self.root):
            return
        
        self.attached_files.remove(file_name)
        if file_name in self.file_paths:
            del self.file_paths[file_name]
        
        for i in range(self.file_listbox.size()):
            if self.file_listbox.get(i) == file_name:
                self.file_listbox.delete(i)
                break
        
        if file_name in self.thumbnails:
            del self.thumbnails[file_name]
        
        if self.rec_id:
            folder_path = self.get_ncr_folder()
            file_path = os.path.join(folder_path, file_name)
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except Exception as e:
                print(f"Ошибка удаления файла: {e}")
        
        self.render_gallery()
        self.status_var.set(f"🗑️ Файл '{file_name}' удален")
    
    def on_file_select(self, event):
        """Обработка выбора файла из списка"""
        selection = self.file_listbox.curselection()
        if not selection:
            return
        
        index = selection[0]
        file_name = self.file_listbox.get(index)
        self.selected_file = file_name
        self.status_var.set(f"📄 Выбран файл: {file_name}")
    
    def on_file_double_click(self, event):
        """Обработка двойного клика по файлу"""
        selection = self.file_listbox.curselection()
        if not selection:
            return
        
        index = selection[0]
        file_name = self.file_listbox.get(index)
        file_path = self.file_paths.get(file_name)
        self.open_file(file_path)
    
    def remove_selected_file(self):
        """Удаляет выбранный файл из списка"""
        selection = self.file_listbox.curselection()
        if not selection:
            messagebox.showwarning("Внимание", "Выберите файл для удаления", parent=self.root)
            return
        
        index = selection[0]
        file_name = self.file_listbox.get(index)
        self.remove_file(file_name)
    
    def open_files_folder(self):
        """Открывает папку с файлами рекламации"""
        if not self.rec_id:
            messagebox.showinfo("Информация", "Сначала сохраните рекламацию", parent=self.root)
            return
        
        folder_path = self.get_ncr_folder()
        if os.path.exists(folder_path):
            os.startfile(folder_path)
        else:
            messagebox.showinfo("Информация", "Папка с файлами еще не создана", parent=self.root)
    
    def update_comboboxes(self):
        """Обновляет значения в комбобоксах"""
        if 'model' in self.comboboxes:
            self.comboboxes['model']['values'] = self.models
            if not self.comboboxes['model'].get():
                self.comboboxes['model'].set('')
        
        if 'supplier' in self.comboboxes:
            self.comboboxes['supplier']['values'] = self.suppliers
            if not self.comboboxes['supplier'].get():
                self.comboboxes['supplier'].set('')
        
        if 'commodity' in self.comboboxes:
            self.comboboxes['commodity']['values'] = self.commodities
            if not self.comboboxes['commodity'].get():
                self.comboboxes['commodity'].set('')
        
        if 'parts_disposal' in self.comboboxes:
            self.comboboxes['parts_disposal']['values'] = self.parts_disposal
            if not self.comboboxes['parts_disposal'].get():
                self.comboboxes['parts_disposal'].set('')
        
        if 'repetition' in self.comboboxes:
            self.comboboxes['repetition']['values'] = self.repetition_list
            if not self.comboboxes['repetition'].get():
                self.comboboxes['repetition'].set('')
        
        if 'defect' in self.comboboxes:
            self.comboboxes['defect']['values'] = self.defects_list
            if not self.comboboxes['defect'].get():
                self.comboboxes['defect'].set('')
    
    def refresh_reference_data(self):
        """Обновляет справочные данные из базы"""
        self.load_reference_data()
        self.update_comboboxes()
        self.status_var.set("✅ Справочники обновлены")
    
    def setup_bindings(self):
        """Настройка горячих клавиш"""
        self.root.bind('<Control-s>', lambda e: self.save_reclamation())
        self.root.bind('<Control-Shift-C>', lambda e: self.clear_fields())
        self.root.bind('<F5>', lambda e: self.refresh_reference_data())
        self.root.bind('<Escape>', lambda e: self.root.quit())
    
    def save_reclamation(self):
        """Сохраняет или обновляет рекламацию в базе данных"""
        self.status_var.set("⏳ Сохранение...")
        self.root.update()
        
        try:
            data = {}
            for key, var in self.vars.items():
                if key == 'report8d_checkbox':
                    data[key] = var.get()
                else:
                    data[key] = var.get().strip()

            if data.get('vin') and len(data['vin']) > 17:
                messagebox.showwarning("Ошибка", "VIN не может быть длиннее 17 символов!", parent=self.root)
                self.focus_window()
                return
            
            if data.get('failure_quantity'):
                try:
                    data['failure_quantity'] = int(data['failure_quantity'])
                except ValueError:
                    messagebox.showwarning("Внимание", "Поле 'Количество отказов' должно быть числом", parent=self.root)
                    self.status_var.set("❌ Ошибка: неверный формат количества отказов")
                    self.focus_window()
                    return
            else:
                data['failure_quantity'] = None
            
            if data.get('cutoff'):
                try:
                    data['cutoff'] = datetime.strptime(data['cutoff'], '%Y-%m-%d')
                except ValueError:
                    messagebox.showwarning("Внимание", "Неверный формат даты cutoff. Используйте ГГГГ-ММ-ДД", parent=self.root)
                    self.status_var.set("❌ Ошибка: неверный формат даты")
                    self.focus_window()
                    return
            else:
                data['cutoff'] = None
            
            data['description'] = self.description_text.get('1.0', 'end-1c').strip()
            data['comments'] = self.comments_text.get('1.0', 'end-1c').strip()
            data['date_creation'] = datetime.now()
            
            required_fields = ['model', 'partnumber', 'supplier']
            missing = [f for f in required_fields if not data.get(f)]
            
            if missing:
                field_names = {
                    'model': 'Модель',
                    'partnumber': 'Номер детали',
                    'supplier': 'Поставщик'
                }
                msg = "Заполните обязательные поля:\n" + "\n".join(f"• {field_names[f]}" for f in missing)
                messagebox.showwarning("Внимание", msg, parent=self.root)
                self.status_var.set("❌ Ошибка: не заполнены обязательные поля")
                self.focus_window()
                return
            
            if self.edit_mode:
                result = self.db.update_reclamation(self.rec_id, data)
                if result['success']:
                    self.save_files_after_save()
                    messagebox.showinfo("Успех", f"✅ Рекламация #{self.rec_id} обновлена!", parent=self.root)
                    self.status_var.set(f"✅ Рекламация #{self.rec_id} обновлена")
                    self.root.after(100, self.focus_window)
                else:
                    messagebox.showerror("Ошибка", f"❌ {result['message']}", parent=self.root)
                    self.status_var.set(f"❌ Ошибка: {result['message']}")
                    self.root.after(100, self.focus_window)
            else:
                result = self.db.save_reclamation(data)
                
                if result['success']:
                    self.rec_id = result['id']
                    current_model = self.vars['model'].get().strip()
                    full_pir = f"AGS-{current_model}-{self.rec_id:04d}"
                    self.db.update_reclamation(self.rec_id, {'full_pir_number': full_pir})
                    self.save_files_after_save()
                    
                    messagebox.showinfo("Успех", f"✅ Рекламация сохранена!\nID: {self.rec_id}\nНомер: {full_pir}", parent=self.root)
                    self.status_var.set(f"✅ Рекламация #{self.rec_id} сохранена")
                    
                    self.edit_mode = True
                    self.root.title(f"Редактирование NCR #{self.rec_id}")
                    
                    self.clear_fields()
                    self.refresh_reference_data()
                    
                    self.status_var.set(f"✅ Рекламация #{self.rec_id} сохранена. Можете создать новую или редактировать текущую.")
                    self.root.after(100, self.focus_window)
                else:
                    messagebox.showerror("Ошибка", f"❌ {result['message']}", parent=self.root)
                    self.status_var.set(f"❌ Ошибка: {result['message']}")
                    self.root.after(100, self.focus_window)
            
        except Exception as e:
            messagebox.showerror("Ошибка", f"❌ Ошибка при сохранении:\n{str(e)}", parent=self.root)
            self.status_var.set(f"❌ Ошибка: {str(e)}")
            self.root.after(100, self.focus_window)
    
    def focus_window(self):
        """Возвращает фокус на текущее окно"""
        try:
            if self.root.winfo_exists():
                self.root.focus_force()
                self.root.lift()
                self.root.after(50, lambda: self.root.focus_force() if self.root.winfo_exists() else None)
        except Exception as e:
            print(f"Ошибка при возврате фокуса: {e}")
    
    def clear_fields(self):
        """Очищает все поля"""
        for key, var in self.vars.items():
            if key == 'report8d_checkbox':
                var.set(False)
            else:
                var.set('')
        self.description_text.delete('1.0', 'end')
        self.comments_text.delete('1.0', 'end')
        self.attached_files = []
        self.file_paths = {}
        self.file_listbox.delete(0, tk.END)
        self.thumbnails = {}
        self.thumbnail_images = []
        self.render_gallery()
        self.status_var.set("Поля очищены")
    
    def export_to_pptx(self):
        """Экспортирует данные рекламации в PowerPoint по шаблону"""
        if not self.rec_id:
            messagebox.showwarning("Внимание", "Сначала сохраните рекламацию!", parent=self.root)
            return
        
        try:
            from pptx import Presentation
            
            # Путь к шаблону
            template_path = r"X:\SQE_DB\Templates\sqe_template.pptx"
            
            # Проверяем существование шаблона
            if not os.path.exists(template_path):
                messagebox.showerror(
                    "Ошибка", 
                    f"Шаблон не найден по пути:\n{template_path}\n\n"
                    "Пожалуйста, проверьте путь к шаблону.",
                    parent=self.root
                )
                return
            
            # Загружаем данные рекламации
            rec = self.db.get_reclamation_by_id(self.rec_id)
            if not rec:
                messagebox.showerror("Ошибка", f"Рекламация #{self.rec_id} не найдена", parent=self.root)
                return
            
            # Загружаем шаблон
            prs = Presentation(template_path)
            
            # Обходим все слайды и заменяем заполнители
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape.text = self._replace_placeholders(shape.text, rec)
            
            # Путь для сохранения
            output_dir = r"X:\SQE_DB\Reports"
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(output_dir, f"NCR_{rec.id}_{timestamp}.pptx")
            
            # Сохраняем презентацию
            prs.save(output_path)
            
            messagebox.showinfo(
                "Успех", 
                f"✅ Презентация создана!\n\n"
                f"Файл: {output_path}",
                parent=self.root
            )
            
            # Предлагаем открыть файл
            if messagebox.askyesno("Открыть файл", "Открыть созданную презентацию?", parent=self.root):
                os.startfile(output_path)
            
        except ImportError:
            messagebox.showerror(
                "Ошибка", 
                "Библиотека python-pptx не установлена.\n\n"
                "Установите ее командой:\n"
                "pip install python-pptx",
                parent=self.root
            )
        except Exception as e:
            messagebox.showerror("Ошибка", f"Ошибка при создании презентации:\n{str(e)}", parent=self.root)

    def _replace_placeholders(self, text, rec):
        """Заменяет заполнители в тексте на данные из рекламации"""
        replacements = {
            '{{id}}': str(rec.id),
            '{{model}}': rec.model or '',
            '{{commodity}}': rec.commodity or '',
            '{{creator}}': rec.creator or '',
            '{{partnumber}}': rec.partnumber or '',
            '{{partname}}': rec.partname or '',
            '{{supplier}}': rec.supplier or '',
            '{{full_pir_number}}': rec.full_pir_number or '',
            '{{repetition}}': rec.repetition or '',
            '{{ncr_status}}': rec.ncr_status or '',
            '{{failure_quantity}}': str(rec.failure_quantity) if rec.failure_quantity else '',
            '{{vin}}': rec.vin or '',
            '{{defect}}': rec.defect or '',
            '{{parts_disposal}}': rec.parts_disposal or '',
            '{{description}}': rec.description or '',
            '{{comments}}': rec.comments or '',
            '{{date_creation}}': rec.date_creation.strftime('%Y-%m-%d %H:%M') if rec.date_creation else '',
            '{{cutoff}}': rec.cutoff.strftime('%Y-%m-%d') if rec.cutoff else '',
            '{{report8d}}': '✅ Да' if rec.report8d_checkbox else '❌ Нет',
        }
        
        for placeholder, value in replacements.items():
            text = text.replace(placeholder, value)
        
        return text